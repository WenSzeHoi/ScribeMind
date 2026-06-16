import os
import tempfile
import urllib.request
import urllib.error
from datetime import datetime
from pathlib import Path

from langchain.tools import tool
from ddgs import DDGS
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH


@tool
def web_search(query: str) -> str:
    """Search the web for current information. Use this for recent events, news, or facts you're unsure about."""
    try:
        results = DDGS().text(query, max_results=10)
        if not results:
            return "No results found."
        return "\n\n".join(
            f"Title: {r['title']}\nURL: {r['href']}\nSnippet: {r['body']}"
            for r in results
        )
    except Exception as e:
        return f"Search error: {e}"


@tool
def search_images(query: str, max_results: int = 5) -> str:
    """Search the web for images related to a topic.

    Use this tool to find relevant pictures, photos, or diagrams that can be
    embedded into a PDF report. Returns image URLs with metadata.

    Args:
        query: What kind of images to search for (e.g. 'Eiffel Tower at night').
        max_results: Maximum number of image results to return (default 5).
    """
    try:
        results = DDGS().images(query, max_results=max_results)
        if not results:
            return "No image results found."

        output_lines = []
        for i, r in enumerate(results, 1):
            title = r.get("title", "Untitled")
            url = r.get("image", "")
            source = r.get("url", "")
            width = r.get("width", "?")
            height = r.get("height", "?")
            output_lines.append(
                f"{i}. Title: {title}\n"
                f"   Image URL: {url}\n"
                f"   Source: {source}\n"
                f"   Dimensions: {width}x{height}"
            )
        return "\n\n".join(output_lines)
    except Exception as e:
        return f"Image search error: {e}"


class ReportPDF(FPDF):
    """Custom PDF with header and footer for reports."""

    def header(self):
        self.set_font("Helvetica", "B", 12)
        self.cell(0, 10, "Research Report", align="C", new_x="LMARGIN", new_y="NEXT")
        self.line(self.l_margin, self.get_y(), self.w - self.r_margin, self.get_y())
        self.ln(4)

    def footer(self):
        self.set_y(-15)
        self.set_font("Helvetica", "I", 8)
        self.cell(0, 10, f"Page {self.page_no()}/{{nb}}", align="C")


def _sanitize(text: str) -> str:
    """Replace characters unsupported by latin-1 (Helvetica) with safe equivalents."""
    result = []
    for ch in text:
        try:
            ch.encode("latin-1")
            result.append(ch)
        except UnicodeEncodeError:
            # Replace common Unicode chars with ASCII equivalents
            replacements = {
                "‘": "'", "’": "'", "“": '"', "”": '"',
                "–": "--", "—": "---", "•": "-",
                "…": "...", " ": " ",
            }
            result.append(replacements.get(ch, "?"))
    return "".join(result)


def _download_image(url: str, dest_dir: str) -> str | None:
    """Download an image from a URL to a local directory. Returns the local path or None."""
    try:
        # Try to get a reasonable file extension from the URL
        path = Path(url.split("?")[0])  # strip query params
        suffix = path.suffix.lower()
        if suffix not in (".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp"):
            suffix = ".jpg"  # default

        dest = os.path.join(dest_dir, f"img_{hash(url) & 0xFFFF:04x}{suffix}")

        req = urllib.request.Request(
            url,
            headers={"User-Agent": "Mozilla/5.0 (compatible; ScribeMind/1.0)"},
        )
        with urllib.request.urlopen(req, timeout=15) as resp:
            with open(dest, "wb") as f:
                f.write(resp.read())
        return dest
    except Exception:
        return None


@tool
def generate_pdf_report(title: str, content: str, images: str = "") -> str:
    """Generate a PDF report and save it to the reports/ directory.

    Use this tool when the user asks you to create a report, write up findings,
    or save search results as a document. Always use after gathering information
    via web_search and optionally search_images for relevant pictures.

    Args:
        title: A concise report title (used as filename and document heading).
        content: The full report body. Use markdown-style formatting:
                 - Use '# Section' for major section headings
                 - Use '## Subsection' for sub-headings
                 - Use '- bullet' for bullet points
                 - Use '[IMG]' as a placeholder where you want an image inserted
                 - Include source URLs inline as plain text
        images: Comma-separated list of image URLs to embed (from search_images).
                Images are placed at [IMG] markers in order, or at the end if no
                markers are present. Example: 'https://example.com/a.jpg, https://example.com/b.png'
    """
    try:
        os.makedirs("reports", exist_ok=True)

        # Parse image URLs
        image_urls = [u.strip() for u in images.split(",") if u.strip()] if images else []

        timestamp = datetime.now().strftime("%Y-%m-%d_%H%M%S")
        safe_title = "".join(c if c.isalnum() or c in " _-" else "_" for c in title)[:80]
        filename = f"reports/{safe_title}_{timestamp}.pdf"

        # Download images to a temp directory
        tmpdir = tempfile.mkdtemp(prefix="scribemind_img_")
        local_images = []
        for url in image_urls:
            local_path = _download_image(url, tmpdir)
            if local_path:
                local_images.append(local_path)

        pdf = ReportPDF()
        pdf.alias_nb_pages()
        pdf.set_auto_page_break(auto=True, margin=20)
        pdf.add_page()

        # Title
        pdf.set_font("Helvetica", "B", 20)
        pdf.multi_cell(0, 10, _sanitize(title), align="L")
        pdf.ln(2)

        # Metadata line
        pdf.set_font("Helvetica", "I", 9)
        pdf.set_text_color(100, 100, 100)
        pdf.cell(0, 6, f"Generated: {datetime.now().strftime('%B %d, %Y at %H:%M')}", new_x="LMARGIN", new_y="NEXT")
        pdf.ln(6)
        pdf.set_text_color(0, 0, 0)

        # Content — parse sections, bullets, paragraphs, and [IMG] markers
        img_index = 0
        for line in content.split("\n"):
            stripped = line.strip()
            if not stripped:
                pdf.ln(4)
                continue

            if stripped == "[IMG]":
                # Insert the next available image
                if img_index < len(local_images):
                    _place_image(pdf, local_images[img_index], f"Image {img_index + 1}")
                    img_index += 1
                else:
                    pdf.set_font("Helvetica", "I", 9)
                    pdf.set_text_color(150, 150, 150)
                    pdf.cell(0, 6, "[Image placeholder — no image available]", new_x="LMARGIN", new_y="NEXT")
                    pdf.set_text_color(0, 0, 0)
                    pdf.ln(4)
                continue

            if stripped.startswith("# "):
                pdf.set_font("Helvetica", "B", 14)
                pdf.ln(4)
                pdf.multi_cell(0, 8, _sanitize(stripped[2:]), align="L")
                pdf.ln(2)
            elif stripped.startswith("## "):
                pdf.set_font("Helvetica", "B", 12)
                pdf.ln(3)
                pdf.multi_cell(0, 7, _sanitize(stripped[3:]), align="L")
                pdf.ln(1)
            elif stripped.startswith("- ") or stripped.startswith("* "):
                pdf.set_font("Helvetica", "", 10)
                pdf.set_x(pdf.l_margin + 8)
                pdf.multi_cell(pdf.w - pdf.l_margin - pdf.r_margin - 8, 6, "- " + _sanitize(stripped[2:]), align="L")
            else:
                pdf.set_font("Helvetica", "", 10)
                pdf.multi_cell(0, 6, _sanitize(stripped), align="L")

        # Any remaining images not placed via [IMG] markers — append at the end
        if img_index < len(local_images):
            pdf.ln(6)
            pdf.set_font("Helvetica", "B", 12)
            pdf.cell(0, 8, "Related Images", new_x="LMARGIN", new_y="NEXT")
            pdf.ln(6)
            while img_index < len(local_images):
                _place_image(pdf, local_images[img_index], f"Image {img_index + 1}")
                img_index += 1

        pdf.output(filename)

        # Clean up temp images
        for f in os.listdir(tmpdir):
            os.unlink(os.path.join(tmpdir, f))
        os.rmdir(tmpdir)

        return f"PDF report saved to: {os.path.abspath(filename)}"
    except Exception as e:
        return f"PDF generation error: {e}"


@tool
def generate_ppt_report(title: str, content: str, images: str = "") -> str:
    """Generate a PowerPoint report and save it to the reports/ directory.

    Use this tool when the user asks you to create a presentation, slides,
    or a PowerPoint deck based on search findings. Always use after gathering
    information via web_search and optionally search_images for relevant pictures.

    Args:
        title: A concise report title (used as filename and slide deck heading).
        content: The full report body. Use markdown-style formatting:
                 - Use '# Section' for major section headings (creates a section
                   header slide)
                 - Use '## Subsection' for sub-headings (creates a content slide
                   with that title)
                 - Use '- bullet' or '* bullet' for bullet points
                 - Use '[IMG]' as a placeholder where you want an image inserted
                 - Include source URLs inline as plain text
        images: Comma-separated list of image URLs to embed (from search_images).
                Images are placed at [IMG] markers in order, or on separate slides
                at the end if no markers are present.
                Example: 'https://example.com/a.jpg, https://example.com/b.png'
    """
    try:
        os.makedirs("reports", exist_ok=True)

        # Parse image URLs
        image_urls = [u.strip() for u in images.split(",") if u.strip()] if images else []

        timestamp = datetime.now().strftime("%Y-%m-%d_%H%M%S")
        safe_title = "".join(c if c.isalnum() or c in " _-" else "_" for c in title)[:80]
        filename = f"reports/{safe_title}_{timestamp}.pptx"

        # Download images to a temp directory
        tmpdir = tempfile.mkdtemp(prefix="scribemind_ppt_img_")
        local_images = []
        for url in image_urls:
            local_path = _download_image(url, tmpdir)
            if local_path:
                local_images.append(local_path)

        prs = Presentation()

        # --- Title slide (layout 0) ---
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        subtitle = slide.shapes.placeholders[1]
        subtitle.text = f"Generated: {datetime.now().strftime('%B %d, %Y at %H:%M')}"

        img_index = 0
        current_slide = None  # tracks the active content slide for body text

        for line in content.split("\n"):
            stripped = line.strip()
            if not stripped:
                continue

            # [IMG] marker: place next image on current slide
            if stripped == "[IMG]":
                if current_slide is None:
                    blank_layout = prs.slide_layouts[6]  # Blank
                    current_slide = prs.slides.add_slide(blank_layout)
                if img_index < len(local_images):
                    _place_image_on_slide(prs, current_slide, local_images[img_index])
                    img_index += 1
                continue

            # '# Section' -> Section Header slide (layout 2)
            if stripped.startswith("# "):
                section_layout = prs.slide_layouts[2]
                slide = prs.slides.add_slide(section_layout)
                slide.shapes.title.text = stripped[2:]
                current_slide = None  # section headers are standalone
                continue

            # '## Subsection' -> Title and Content slide (layout 1)
            if stripped.startswith("## "):
                content_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(content_layout)
                slide.shapes.title.text = stripped[3:]
                current_slide = slide
                continue

            # Body text — needs a content slide to write into
            if current_slide is None:
                content_layout = prs.slide_layouts[1]
                current_slide = prs.slides.add_slide(content_layout)
                current_slide.shapes.title.text = ""

            # Access the body text frame on the Title and Content layout
            body_shape = current_slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            # Determine if we should use the first (default) paragraph or add a new one
            if tf.paragraphs[0].text == "" and len(tf.paragraphs) == 1:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_after = Pt(6)

            if stripped.startswith("- ") or stripped.startswith("* "):
                p.text = stripped[2:]
                p.level = 0
                p.font.size = Pt(14)
            else:
                p.text = stripped
                p.font.size = Pt(14)

        # Remaining images not placed via [IMG] markers — each on its own slide
        if img_index < len(local_images):
            blank_layout = prs.slide_layouts[6]
            while img_index < len(local_images):
                slide = prs.slides.add_slide(blank_layout)
                _place_image_on_slide(prs, slide, local_images[img_index])
                img_index += 1

        prs.save(filename)

        # Clean up temp images
        for f in os.listdir(tmpdir):
            os.unlink(os.path.join(tmpdir, f))
        os.rmdir(tmpdir)

        return f"PPT report saved to: {os.path.abspath(filename)}"

    except Exception as e:
        return f"PPT generation error: {e}"


def _place_image(pdf: FPDF, image_path: str, caption: str = "") -> None:
    """Place an image into the PDF, sizing it to fit the page width."""
    try:
        # Get image dimensions from the file
        from PIL import Image
        img = Image.open(image_path)
        img_w, img_h = img.size
        img.close()
    except Exception:
        img_w, img_h = 800, 600  # fallback

    # Calculate display size — max width = page width minus margins, max height = 120mm
    max_w = pdf.w - pdf.l_margin - pdf.r_margin
    max_h = 120

    scale = min(max_w / img_w, max_h / img_h, 1.0)
    disp_w = img_w * scale
    disp_h = img_h * scale

    # Center the image
    x = pdf.l_margin + (max_w - disp_w) / 2

    # Check if we need a page break
    if pdf.get_y() + disp_h + 10 > pdf.h - pdf.b_margin:
        pdf.add_page()

    pdf.ln(2)
    pdf.image(image_path, x=x, w=disp_w, h=disp_h)
    pdf.ln(2)

    if caption:
        pdf.set_font("Helvetica", "I", 8)
        pdf.set_text_color(100, 100, 100)
        pdf.cell(0, 5, _sanitize(caption), align="C", new_x="LMARGIN", new_y="NEXT")
        pdf.set_text_color(0, 0, 0)
    pdf.ln(6)


def _place_image_on_slide(prs, slide, image_path: str) -> None:
    """Place an image onto a PowerPoint slide, centered and scaled to fit."""
    try:
        from PIL import Image
        img = Image.open(image_path)
        img_w, img_h = img.size
        img.close()
    except Exception:
        img_w, img_h = 800, 600  # fallback

    # Convert pixel dimensions to inches (assume 96 DPI, standard for web images)
    dpi = 96
    img_w_in = img_w / dpi
    img_h_in = img_h / dpi

    # Maximum dimensions in inches
    max_w_in = 8.0
    max_h_in = 5.0

    # Scale to fit within bounds
    scale = min(max_w_in / img_w_in, max_h_in / img_h_in, 1.0)
    # Convert to EMUs for python-pptx
    disp_w = Inches(img_w_in * scale)
    disp_h = Inches(img_h_in * scale)

    # Center on the slide
    left = int((prs.slide_width - disp_w) / 2)
    top = int((prs.slide_height - disp_h) / 2)

    slide.shapes.add_picture(image_path, left, top, disp_w, disp_h)


def _place_image_in_docx(doc: Document, image_path: str, caption: str = "") -> None:
    """Place an image into a Word document, centered and scaled to fit the page width."""
    try:
        from PIL import Image
        img = Image.open(image_path)
        img_w, img_h = img.size
        img.close()
    except Exception:
        img_w, img_h = 800, 600  # fallback

    # Calculate display size — max width 6 inches (standard Word margin), max height 5 inches
    dpi = 96
    img_w_in = img_w / dpi
    img_h_in = img_h / dpi

    max_w_in = 6.0
    max_h_in = 5.0

    scale = min(max_w_in / img_w_in, max_h_in / img_h_in, 1.0)
    disp_w = DocxInches(img_w_in * scale)
    disp_h = DocxInches(img_h_in * scale)

    # Add image (centered)
    paragraph = doc.add_paragraph()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = paragraph.add_run()
    run.add_picture(image_path, width=disp_w, height=disp_h)

    if caption:
        cap_para = doc.add_paragraph()
        cap_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        cap_run = cap_para.add_run(caption)
        cap_run.font.size = DocxPt(9)
        cap_run.font.color.rgb = RGBColor(100, 100, 100)
        cap_run.italic = True


@tool
def generate_word_report(title: str, content: str, images: str = "") -> str:
    """Generate a Word document (.docx) and save it to the reports/ directory.

    Use this tool when the user asks you to create a Word document, a .docx file,
    or a report in Microsoft Word format. Always use after gathering information
    via web_search and optionally search_images for relevant pictures.

    Args:
        title: A concise report title (used as filename and document heading).
        content: The full report body. Use markdown-style formatting:
                 - Use '# Section' for major section headings (Heading 1)
                 - Use '## Subsection' for sub-headings (Heading 2)
                 - Use '- bullet' or '* bullet' for bullet points
                 - Use '[IMG]' as a placeholder where you want an image inserted
                 - Include source URLs inline as plain text
        images: Comma-separated list of image URLs to embed (from search_images).
                Images are placed at [IMG] markers in order, or at the end if no
                markers are present. Example: 'https://example.com/a.jpg, https://example.com/b.png'
    """
    try:
        os.makedirs("reports", exist_ok=True)

        # Parse image URLs
        image_urls = [u.strip() for u in images.split(",") if u.strip()] if images else []

        timestamp = datetime.now().strftime("%Y-%m-%d_%H%M%S")
        safe_title = "".join(c if c.isalnum() or c in " _-" else "_" for c in title)[:80]
        filename = f"reports/{safe_title}_{timestamp}.docx"

        # Download images to a temp directory
        tmpdir = tempfile.mkdtemp(prefix="scribemind_docx_img_")
        local_images = []
        for url in image_urls:
            local_path = _download_image(url, tmpdir)
            if local_path:
                local_images.append(local_path)

        doc = Document()

        # --- Default document styles ---
        style = doc.styles["Normal"]
        font = style.font
        font.name = "Calibri"
        font.size = DocxPt(11)

        # --- Title ---
        title_para = doc.add_paragraph()
        title_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
        title_run = title_para.add_run(title)
        title_run.font.size = DocxPt(22)
        title_run.bold = True

        # --- Metadata line ---
        meta_para = doc.add_paragraph()
        meta_run = meta_para.add_run(
            f"Generated: {datetime.now().strftime('%B %d, %Y at %H:%M')}"
        )
        meta_run.font.size = DocxPt(9)
        meta_run.font.color.rgb = RGBColor(100, 100, 100)
        meta_run.italic = True

        # Add a horizontal rule (thin border)
        hr_para = doc.add_paragraph()
        hr_para.paragraph_format.space_before = DocxPt(4)
        hr_para.paragraph_format.space_after = DocxPt(12)
        pPr = hr_para._p.get_or_add_pPr()
        from docx.oxml.ns import qn
        pBdr = pPr.makeelement(qn("w:pBdr"), {})
        bottom = pBdr.makeelement(qn("w:bottom"), {
            qn("w:val"): "single",
            qn("w:sz"): "6",
            qn("w:space"): "1",
            qn("w:color"): "999999",
        })
        pBdr.append(bottom)
        pPr.append(pBdr)

        # --- Parse content ---
        img_index = 0
        for line in content.split("\n"):
            stripped = line.strip()
            if not stripped:
                continue

            if stripped == "[IMG]":
                if img_index < len(local_images):
                    _place_image_in_docx(doc, local_images[img_index], f"Image {img_index + 1}")
                    img_index += 1
                else:
                    placeholder = doc.add_paragraph()
                    placeholder.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    pr = placeholder.add_run("[Image placeholder — no image available]")
                    pr.font.size = DocxPt(9)
                    pr.font.color.rgb = RGBColor(150, 150, 150)
                    pr.italic = True
                continue

            if stripped.startswith("# "):
                heading = doc.add_heading(stripped[2:], level=1)
                heading.paragraph_format.space_before = DocxPt(18)
                continue

            if stripped.startswith("## "):
                heading = doc.add_heading(stripped[3:], level=2)
                heading.paragraph_format.space_before = DocxPt(14)
                continue

            if stripped.startswith("- ") or stripped.startswith("* "):
                para = doc.add_paragraph(stripped[2:], style="List Bullet")
            else:
                para = doc.add_paragraph(stripped)

        # Remaining images not placed via [IMG] markers — append at the end
        if img_index < len(local_images):
            doc.add_heading("Related Images", level=2)
            while img_index < len(local_images):
                _place_image_in_docx(doc, local_images[img_index], f"Image {img_index + 1}")
                img_index += 1

        doc.save(filename)

        # Clean up temp images
        for f in os.listdir(tmpdir):
            os.unlink(os.path.join(tmpdir, f))
        os.rmdir(tmpdir)

        return f"Word report saved to: {os.path.abspath(filename)}"
    except Exception as e:
        return f"Word generation error: {e}"
